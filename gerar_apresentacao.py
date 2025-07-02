
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# --- Funções Auxiliares ---
def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text

def add_section_header_slide(prs, title_text):
    slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text

def add_content_slide(prs, title_text, content_list):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)

def add_image_slide(prs, title_text, img_path, left=Inches(1.5), top=Inches(1.5), width=Inches(7)):
    slide_layout = prs.slide_layouts[5] # Layout "Título Apenas"
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    slide.shapes.add_picture(img_path, left, top, width=width)

# --- Criação da Apresentação ---
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

print("Criando a apresentação...")

# 1. Slide de Título
add_title_slide(prs, 
                "Aplicação de Aprendizado Supervisionado na Indústria Cinematográfica", 
                "Integrantes: Felipe Borges e Fernando Filter")

# 2. Seção: Análise de Sentimentos
add_section_header_slide(prs, "Tarefa 1: Análise de Sentimentos")
add_content_slide(prs, "Análise de Sentimentos: Objetivo e Metodologia",
                  ["Objetivo: Classificar críticas de filmes como positivas ou negativas.",
                   "Dataset: IMDb Large Movie Review Dataset (50.000 críticas).",
                   "Técnica: Vetorização com TF-IDF.",
                   "Modelo: Regressão Logística.",
                   "Resultado: Acurácia de 89.47%."])
add_image_slide(prs, "Distribuição das Classes (Balanceado)", "graficos/sentimentos_distribuicao.png", width=Inches(6), left=Inches(5))
add_image_slide(prs, "Nuvem de Palavras: Críticas Positivas", "graficos/sentimentos_wordcloud_pos.png")
add_image_slide(prs, "Nuvem de Palavras: Críticas Negativas", "graficos/sentimentos_wordcloud_neg.png")

# 3. Seção: Previsão de Bilheteria
add_section_header_slide(prs, "Tarefa 2: Previsão de Sucesso de Bilheteria")
add_content_slide(prs, "Previsão de Bilheteria: Objetivo e Metodologia",
                  ["Objetivo: Prever se um filme terá ROI (Retorno/Orçamento) >= 3.",
                   "Dataset: TMDB 5000 Movie Dataset.",
                   "Features: Orçamento, popularidade, duração, nota média, contagem de votos.",
                   "Modelo: Random Forest Classifier.",
                   "Resultado: Acurácia de 72.45%."])
add_image_slide(prs, "Matriz de Correlação entre Features", "graficos/bilheteria_correlacao.png")
add_image_slide(prs, "Distribuição das Features Numéricas", "graficos/bilheteria_histogramas.png")

# 4. Seção: Classificação de Pôsteres
add_section_header_slide(prs, "Tarefa 3: Classificação de Gênero por Pôster")
add_content_slide(prs, "Classificação de Pôsteres: Objetivo e Metodologia",
                  ["Objetivo: Classificar o gênero de um filme a partir de seu pôster (Multi-label).",
                   "Dataset: Movie Poster Dataset (amostra de 500 pôsteres).",
                   "Técnica: Rede Neural Convolucional (CNN).",
                   "Modelo: Keras/TensorFlow com 3 camadas convolucionais.",
                   "Resultado: Acurácia de 57.53% (prova de conceito)."])
add_image_slide(prs, "Distribuição dos 10 Gêneros Mais Comuns", "graficos/posters_distribuicao_generos.png")

# 5. Slide de Conclusão
add_section_header_slide(prs, "Conclusão")
add_content_slide(prs, "Resultados e Próximos Passos",
                  ["Sucesso na aplicação de ML em 3 domínios: Texto, Tabelas e Imagens.",
                   "Análise de Sentimentos: Modelo robusto com 89.47% de acurácia.",
                   "Previsão de Bilheteria: Modelo funcional com 72.45% de acurácia.",
                   "Classificação de Pôsteres: Pipeline estabelecido com 57.53% de acurácia inicial.",
                   "Próximos Passos: Aumentar datasets, treinar por mais tempo e otimizar hiperparâmetros."])

# Salvar a apresentação
output_filename = "Apresentacao_Final.pptx"
prs.save(output_filename)

print(f"Apresentação salva como '{output_filename}'")
